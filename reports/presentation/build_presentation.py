# Builds reports/cs2_match_prediction_presentation.pptx
#
# 15-slide, 15-minute university presentation for the CS2 match prediction
# project. Every number is transcribed from the frozen reports listed in
# reports/presentation_sources.md — nothing is recomputed here.
#
# Design: light academic theme — white background, navy ink, restrained
# blue/cyan accents (deliberately NOT the dark esports theme of the PWA).

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "reports" / "figures"
PFIG = FIG / "presentation"
OUT = ROOT / "reports" / "cs2_match_prediction_presentation.pptx"

# ---- design tokens ---------------------------------------------------------
NAVY = RGBColor(0x0F, 0x2A, 0x43)   # primary ink
INK2 = RGBColor(0x44, 0x58, 0x6C)   # secondary ink
GRAY = RGBColor(0x74, 0x84, 0x94)   # muted
BLUE = RGBColor(0x1D, 0x6F, 0xB8)   # accent
CYAN = RGBColor(0x35, 0x8C, 0xA8)   # accent 2
ORANGE = RGBColor(0xD9, 0x77, 0x2A) # reality highlight (used sparingly)
GREEN = RGBColor(0x3E, 0x8E, 0x6B)
PANEL = RGBColor(0xF2, 0xF7, 0xFB)  # light panel fill
PANEL2 = RGBColor(0xE9, 0xF1, 0xF8)
LINE = RGBColor(0xC9, 0xD7, 0xE4)   # hairlines
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
FONT_SB = "Segoe UI Semibold"

SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.62)                    # content left margin
CW = Inches(12.09)                   # content width

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

TOTAL = 15


# ---- helpers ---------------------------------------------------------------
def _set_text(tf, lines, default_size=14, default_color=NAVY, default_font=FONT,
              align=PP_ALIGN.LEFT, line_spacing=1.0, space_after=4):
    """lines: list of dicts or list-of-runs. dict keys: runs(list of (text, opts)),
    size, bold, color, font, align, level, space_before/after, spacing."""
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ln.get("align", align)
        p.line_spacing = ln.get("spacing", line_spacing)
        p.space_after = Pt(ln.get("space_after", space_after))
        p.space_before = Pt(ln.get("space_before", 0))
        runs = ln.get("runs")
        if runs is None:
            runs = [(ln.get("text", ""), {})]
        for text, opts in runs:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", ln.get("font", default_font))
            f.size = Pt(opts.get("size", ln.get("size", default_size)))
            f.bold = opts.get("bold", ln.get("bold", False))
            f.italic = opts.get("italic", ln.get("italic", False))
            f.color.rgb = opts.get("color", ln.get("color", default_color))


def tbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, **kw):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _set_text(tf, lines, **kw)
    return box


def panel(slide, x, y, w, h, fill=PANEL, line=LINE, radius=0.06, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    shp.text_frame.margin_left = shp.text_frame.margin_right = Inches(0.12)
    shp.text_frame.margin_top = shp.text_frame.margin_bottom = Inches(0.08)
    return shp


def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def arrow_r(slide, x, y, w, h, fill=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.adjustments[0] = 0.55
    shp.adjustments[1] = 0.55
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def arrow_d(slide, x, y, w, h, fill=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def pic(slide, path, x, y, max_w, max_h, align="center", frame=True):
    """Place image preserving aspect ratio inside (max_w, max_h)."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar >= box_ar:
        w = max_w
        h = int(max_w / ar)
    else:
        h = max_h
        w = int(max_h * ar)
    px = x + (max_w - w) // 2 if align == "center" else x
    py = y + (max_h - h) // 2
    p = slide.shapes.add_picture(str(path), px, py, width=w, height=h)
    if frame:
        p.line.color.rgb = LINE
        p.line.width = Pt(1)
    return p


def header(slide, idx, kicker, title, title_size=26):
    rect(slide, 0, 0, Inches(0.14), SH, BLUE)  # left accent spine
    tbox(slide, ML, Inches(0.30), CW, Inches(0.28),
         [{"runs": [(f"{idx:02d}", {"color": BLUE, "bold": True, "size": 12}),
                    ("   ·   " + kicker.upper(), {"color": GRAY, "size": 12})]}])
    tbox(slide, ML, Inches(0.56), CW, Inches(0.62),
         [{"text": title, "size": title_size, "bold": True, "font": FONT_SB}])
    rect(slide, ML, Inches(1.18), CW, Pt(1.6), LINE)


def footer(slide, idx):
    tbox(slide, ML, Inches(7.14), Inches(8.0), Inches(0.3),
         [{"text": "CS2 Match Prediction — probabilistic modelling & Major simulation",
           "size": 9.5, "color": GRAY}])
    tbox(slide, Inches(11.6), Inches(7.14), Inches(1.1), Inches(0.3),
         [{"text": f"{idx} / {TOTAL}", "size": 9.5, "color": GRAY,
           "align": PP_ALIGN.RIGHT}])


def notes(slide, main, script, understand, qa, transition):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    parts = ["MAIN MESSAGE", main, "", "SCRIPT (speak naturally)", script, "",
             "WHAT I MUST UNDERSTAND"]
    parts += ["- " + u for u in understand]
    parts += ["", "LIKELY PROFESSOR QUESTIONS"]
    for q, a in qa:
        parts.append("Q: " + q)
        parts.append("A: " + a)
    parts += ["", "TRANSITION", transition]
    tf.text = parts[0]
    for line in parts[1:]:
        p = tf.add_paragraph()
        p.text = line


def new_slide():
    return prs.slides.add_slide(BLANK)


def chip(slide, x, y, w, h, number, label, num_color=NAVY, num_size=24, label_size=11):
    panel(slide, x, y, w, h)
    tbox(slide, x + Inches(0.14), y + Inches(0.10), w - Inches(0.28), h - Inches(0.2),
         [{"text": number, "size": num_size, "bold": True, "color": num_color,
           "font": FONT_SB, "space_after": 1},
          {"text": label, "size": label_size, "color": INK2, "spacing": 1.0}])


def bullet_lines(items, size=14, gap=6, marker_color=BLUE):
    out = []
    for it in items:
        if isinstance(it, str):
            it = {"text": it}
        runs = it.get("runs", [(it.get("text", ""), {})])
        out.append({"runs": [("▪  ", {"color": marker_color, "size": size - 1})] + runs,
                    "size": it.get("size", size), "space_after": gap,
                    "spacing": it.get("spacing", 1.04)})
    return out


def styled_table(slide, x, y, w, h, data, col_widths=None, header_fill=NAVY,
                 body_size=12.5, header_size=12.5, row_h=None, align_map=None):
    rows, cols = len(data), len(data[0])
    g = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = g.table
    # kill banding
    tbl = g._element.graphic.graphicData.tbl
    tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # "no style" GUID
    if col_widths:
        for i, cwid in enumerate(col_widths):
            table.columns[i].width = cwid
    if row_h:
        for r in range(rows):
            table.rows[r].height = row_h
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else PANEL
            tf = cell.text_frame
            tf.word_wrap = True
            if isinstance(val, tuple):
                text, opts = val
            else:
                text, opts = val, {}
            p = tf.paragraphs[0]
            if align_map and ci in align_map:
                p.alignment = align_map[ci]
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(opts.get("size", header_size if ri == 0 else body_size))
            f.bold = opts.get("bold", ri == 0)
            f.color.rgb = opts.get("color", WHITE if ri == 0 else NAVY)
    return table


# ===========================================================================
# SLIDE 1 — TITLE
# ===========================================================================
s = new_slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, Inches(0.14), SH, BLUE)
rect(s, Inches(0.62), Inches(2.02), Inches(1.6), Pt(3), BLUE)
tbox(s, Inches(0.62), Inches(1.28), Inches(12.0), Inches(0.5),
     [{"text": "MACHINE LEARNING PROJECT PRESENTATION", "size": 13, "color": GRAY}])
tbox(s, Inches(0.62), Inches(2.22), Inches(12.1), Inches(1.9),
     [{"text": "Predicting Counter-Strike 2 Matches", "size": 40, "bold": True,
       "font": FONT_SB, "space_after": 2},
      {"runs": [("with ", {"size": 40, "bold": True, "font": FONT_SB}),
                ("Probabilities", {"size": 40, "bold": True, "font": FONT_SB, "color": BLUE})],
       "space_after": 10},
      {"text": "From pre-match data to a fully simulated Major — and back to reality",
       "size": 17, "color": INK2}])
cx = Inches(0.62)
cw4 = Inches(2.85)
gap = Inches(0.23)
labels = [("9,923", "matches, Jan 2023 → Jun 2026"),
          ("3", "model families compared (LR · RF · XGB)"),
          ("50,000", "Monte-Carlo simulated Majors"),
          ("1", "deployed prediction app (API + PWA)")]
for i, (n, l) in enumerate(labels):
    chip(s, cx + i * (cw4 + gap), Inches(4.55), cw4, Inches(1.05), n, l, num_color=BLUE)
tbox(s, Inches(0.62), Inches(6.30), Inches(11.0), Inches(0.6),
     [{"text": "Yiğit Alp Öztorun   ·   August 2026", "size": 14, "color": NAVY, "bold": True}])
notes(s,
      "This project builds and honestly evaluates a probabilistic prediction system for professional Counter-Strike 2 matches.",
      "Good morning. In this presentation I'll show a machine-learning system that predicts professional Counter-Strike 2 matches. "
      "The key word on this slide is 'probabilities' — the system never just says 'Team A will win'; it estimates how likely each team is to win. "
      "I'll walk through the data, the methodology, three model families that were compared, the two final frozen models, and then the part I'm most proud of: "
      "before the IEM Cologne Major 2026 started, the system was frozen and used to simulate the whole tournament 50,000 times — "
      "and afterwards we compared those frozen predictions against what actually happened. Everything ends in a deployed web application.",
      ["CS2 = Counter-Strike 2, a 5-vs-5 esport; matches are best-of-1/3/5 series of 'maps'.",
       "The whole talk has one arc: data → models → frozen evaluation → real tournament → app."],
      [("Why esports?", "Rich public match history, frequent events, and a genuinely hard, high-variance prediction problem — ideal for studying probabilistic ML honestly."),
       ("Is this betting-related?", "No — it is an academic exercise in probabilistic prediction and honest evaluation.")],
      "Let's define the actual prediction problem precisely.")
footer(s, 1)

# ===========================================================================
# SLIDE 2 — OBJECTIVE
# ===========================================================================
s = new_slide()
header(s, 2, "Objective", "Estimate P(Team A wins) — before the match starts")
tbox(s, ML, Inches(1.45), Inches(6.7), Inches(0.9),
     [{"runs": [("Input, strictly pre-match:  ", {"bold": True, "size": 15}),
                ("the two teams, the series format (BO1/BO3/BO5) and everything both teams did before this match. "
                 "Never a score, never a stat from the match itself.", {"size": 15, "color": INK2})],
       "spacing": 1.12}])
p1 = panel(s, ML, Inches(2.55), Inches(6.7), Inches(1.5))
tbox(s, ML + Inches(0.25), Inches(2.75), Inches(6.2), Inches(1.2),
     [{"text": "Output", "size": 12, "color": GRAY, "space_after": 4},
      {"runs": [("P(Team A beats Team B)", {"size": 24, "bold": True, "color": BLUE, "font": FONT_SB}),
                ("   e.g. 0.62", {"size": 18, "color": INK2})]}])
tbox(s, ML, Inches(4.35), Inches(6.7), Inches(2.6),
     bullet_lines([
         {"runs": [("A label hides information ", {"bold": True}),
                   ("— “A wins” at 51% and at 95% are very different claims.", {"color": INK2})]},
         {"runs": [("Probabilities can be scored ", {"bold": True}),
                   ("— Log Loss and Brier punish overconfident mistakes; accuracy cannot see them.", {"color": INK2})]},
         {"runs": [("Simulation needs probabilities ", {"bold": True}),
                   ("— a tournament simulator must sample upsets, not replay favorites.", {"color": INK2})]},
         {"runs": [("Honest about variance ", {"bold": True}),
                   ("— upsets are a structural feature of CS2; the model should say so.", {"color": INK2})]},
     ], size=14.5, gap=9))
# right visual: classifier vs probability
rx = Inches(7.75)
rw = Inches(4.95)
panel(s, rx, Inches(1.5), rw, Inches(2.28), fill=PANEL)
tbox(s, rx + Inches(0.25), Inches(1.68), rw - Inches(0.5), Inches(0.4),
     [{"text": "PLAIN CLASSIFIER", "size": 11, "color": GRAY, "space_after": 2}])
tbox(s, rx + Inches(0.25), Inches(2.05), rw - Inches(0.5), Inches(0.6),
     [{"text": "“Team A wins.”", "size": 20, "bold": True}])
tbox(s, rx + Inches(0.25), Inches(2.75), rw - Inches(0.5), Inches(0.9),
     [{"text": "One bit of output. No way to say how sure — or to be told it was badly overconfident.",
       "size": 12.5, "color": INK2, "spacing": 1.1}])
panel(s, rx, Inches(4.0), rw, Inches(2.75), fill=WHITE, line=BLUE)
tbox(s, rx + Inches(0.25), Inches(4.18), rw - Inches(0.5), Inches(0.4),
     [{"text": "THIS SYSTEM", "size": 11, "color": BLUE, "space_after": 2}])
tbox(s, rx + Inches(0.25), Inches(4.52), rw - Inches(0.5), Inches(0.5),
     [{"runs": [("Team A  62%", {"size": 19, "bold": True, "color": BLUE}),
                ("     Team B  38%", {"size": 19, "bold": True, "color": INK2})]}])
bar_x = rx + Inches(0.25)
bar_w = rw - Inches(0.5)
rect(s, bar_x, Inches(5.15), bar_w, Inches(0.30), PANEL2, line=LINE)
rect(s, bar_x, Inches(5.15), Emu(int(bar_w * 0.62)), Inches(0.30), BLUE)
tbox(s, rx + Inches(0.25), Inches(5.62), rw - Inches(0.5), Inches(1.0),
     [{"text": "A full distribution over outcomes — scoreable, comparable, and usable as the engine of a Monte-Carlo tournament simulation.",
       "size": 12.5, "color": INK2, "spacing": 1.12}])
notes(s,
      "The task is probability estimation from strictly pre-match information — not just winner classification.",
      "The problem statement is deliberately strict: given only information available before a match starts — who is playing, the series format, and both teams' "
      "history — estimate the probability that Team A wins. Why insist on a probability instead of a predicted winner? Three reasons. "
      "First, a label hides confidence: 'A wins' at 51 percent and at 95 percent are completely different statements. Second, probabilities can be properly scored — "
      "Log Loss and Brier score punish a model that is confidently wrong, which accuracy can't detect. Third, the downstream goal — simulating a whole tournament — "
      "only works with probabilities, because you need to sample upsets with realistic frequency. In an esport where upsets are common, saying '62 percent' honestly "
      "is more useful than pretending certainty.",
      ["'Pre-match only' is the core constraint that drives all leakage prevention later.",
       "Log Loss = −log of the probability assigned to what actually happened, averaged; Brier = mean squared error of the probability.",
       "Both are proper scoring rules: the best long-run strategy is reporting your true belief."],
      [("Why not just maximize accuracy?", "Accuracy only checks which side of 0.5 you were on; it treats 51% and 99% identically. For simulation and honest uncertainty we need probability quality, measured by Log Loss/Brier."),
       ("What makes this hard?", "High variance (upsets), rosters change, the meta shifts over time, and public data has identity/label quality problems we had to fix first.")],
      "Here is the system that answers this question — it actually has two prediction modes.")
footer(s, 2)

# ===========================================================================
# SLIDE 3 — SYSTEM AT A GLANCE
# ===========================================================================
s = new_slide()
header(s, 3, "System overview", "One question, two prediction modes")
# Mode A row
ya = Inches(1.55)
row_h = Inches(2.15)
panel(s, ML, ya, CW, row_h, fill=PANEL)
tbox(s, ML + Inches(0.22), ya + Inches(0.14), Inches(4.2), Inches(0.6),
     [{"runs": [("MODE A — MAPS UNKNOWN", {"size": 12.5, "bold": True, "color": BLUE})],
       "space_after": 1},
      {"text": "pre-veto: before map picks exist", "size": 11, "color": GRAY}])


def flow(slide, y, boxes, box_w, box_h, x0, gap_w=Inches(0.42), fills=None, sizes=None):
    x = x0
    n = len(boxes)
    for i, txt in enumerate(boxes):
        f = (fills[i] if fills else WHITE)
        pn = panel(slide, x, y, box_w, box_h, fill=f, line=LINE)
        tf = pn.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lines = txt if isinstance(txt, list) else [txt]
        _set_text(tf, lines, align=PP_ALIGN.CENTER)
        x += box_w
        if i < n - 1:
            arrow_r(slide, x + Inches(0.05), y + box_h / 2 - Inches(0.09),
                    gap_w - Inches(0.10), Inches(0.18), fill=LINE)
            x += gap_w
    return x


bw = Inches(2.55)
bh = Inches(1.05)
flow(s, ya + Inches(0.80),
     [[{"text": "Team A · Team B · BO1/3/5", "size": 12.5, "bold": True}],
      [{"text": "Random Forest V2", "size": 13, "bold": True, "color": BLUE, "space_after": 1},
       {"text": "17 pre-match series features", "size": 10.5, "color": INK2}],
      [{"text": "P(Team A wins the series)", "size": 12.5, "bold": True}]],
     bw, bh, ML + Inches(1.75), fills=[WHITE, WHITE, PANEL2])
tbox(s, ML + Inches(10.35), ya + Inches(0.86), Inches(1.6), Inches(1.0),
     [{"text": "used by the\ntournament\nsimulator", "size": 10.5, "color": GRAY, "spacing": 1.05}])

# Mode B row
yb = Inches(3.92)
panel(s, ML, yb, CW, row_h, fill=PANEL)
tbox(s, ML + Inches(0.22), yb + Inches(0.14), Inches(5.2), Inches(0.6),
     [{"runs": [("MODE B — MAPS KNOWN", {"size": 12.5, "bold": True, "color": CYAN})],
       "space_after": 1},
      {"text": "after the veto: the ordered maps are input", "size": 11, "color": GRAY}])
bw2 = Inches(2.42)
flow(s, yb + Inches(0.80),
     [[{"text": "Teams + ordered maps", "size": 12.5, "bold": True}],
      [{"text": "XGBoost V3 — per map", "size": 12.5, "bold": True, "color": CYAN, "space_after": 1},
       {"text": "131 map-level features", "size": 10.5, "color": INK2}],
      [{"text": "p₁, p₂, p₃ …", "size": 13, "bold": True, "space_after": 1},
       {"text": "one probability per map", "size": 10.5, "color": INK2}],
      [{"text": "Dynamic program", "size": 12.5, "bold": True, "space_after": 1},
       {"text": "→ P(series win)", "size": 11.5, "color": INK2}]],
     bw2, bh, ML + Inches(1.75), gap_w=Inches(0.34),
     fills=[WHITE, WHITE, WHITE, PANEL2])
tbox(s, ML, Inches(6.35), CW, Inches(0.6),
     [{"runs": [("During development, ", {"color": INK2, "size": 13.5}),
                ("Logistic Regression, Random Forest and XGBoost", {"bold": True, "size": 13.5}),
                (" were tuned and compared on both tasks — the two frozen winners above serve two different application tasks and are never blended.",
                 {"color": INK2, "size": 13.5})], "spacing": 1.12}])
notes(s,
      "There are two frozen models because there are two genuinely different prediction tasks — pre-veto series prediction and known-map prediction.",
      "The system answers the same question in two situations. Mode A: before the map veto, all we know is who plays and the format — a tuned Random Forest, "
      "using 17 pre-match team-history features, directly outputs the series win probability. This is also the engine of the tournament simulator, because for "
      "future tournament matches maps are never known in advance. Mode B: once the map veto has happened, the ordered maps are legitimate input. There, a tuned "
      "XGBoost model predicts each individual map with 131 richer features, and a small dynamic program combines the per-map probabilities into a series probability "
      "— it accounts for the fact that map 3 is only played if the series is still alive. Three model families were compared during development; these two are the "
      "frozen winners for their respective tasks.",
      ["Veto = the pick/ban phase where teams alternately eliminate and pick maps before a series.",
       "RF V2 and XGB V3 solve different tasks (series vs single map) — their accuracies are not comparable numbers.",
       "The DP is exact composition, not simulation."],
      [("Why not one model for both?", "The information sets differ: knowing the exact maps adds real signal (map-specific strength, per-map rosters). A single model would either waste that input or need fake placeholder maps."),
       ("Why RF for one task and XGB for the other?", "Each was selected by frozen, pre-registered criteria on its own task's validation — probability quality first. RF V2 won the series task, XGB V3 the map task.")],
      "Both models learn from the same underlying dataset — let's look at it.")
footer(s, 3)

# ===========================================================================
# SLIDE 4 — DATASET
# ===========================================================================
s = new_slide()
header(s, 4, "Data", "3.5 years of professional CS2 matches")
cw3 = Inches(2.90)
gap = Inches(0.17)
tiles = [("9,923", "series (raw), Jan 2023 → 28 Jun 2026"),
         ("10,674", "individual map results"),
         ("344", "tournaments across 3 tiers"),
         ("BO3", "dominant format — 7,934 BO3 · 1,784 BO1 · 89 BO5")]
for i, (n, l) in enumerate(tiles):
    chip(s, ML + i * (cw3 + gap), Inches(1.5), cw3, Inches(1.12), n, l, num_color=BLUE)
tbox(s, ML, Inches(3.0), Inches(6.6), Inches(3.6),
     bullet_lines([
         {"runs": [("Target:  ", {"bold": True}),
                   ("team1 wins the series — reconstructed from the final score.", {"color": INK2})]},
         {"runs": [("The provided winner label was broken ", {"bold": True}),
                   ("— it disagreed with the actual scores in 49.9% of rows; the audit caught this before any model saw it.", {"color": INK2})]},
         {"runs": [("Team IDs were not persistent ", {"bold": True}),
                   ("— every ID appears in exactly one match; identity had to be rebuilt from normalized team names + manual review.", {"color": INK2})]},
         {"runs": [("Chronological by nature ", {"bold": True}),
                   ("— teams evolve, rosters change; every split and feature respects time order.", {"color": INK2})]},
     ], size=14, gap=10))
p = panel(s, Inches(7.5), Inches(3.0), Inches(5.2), Inches(3.6), fill=PANEL)
tbox(s, Inches(7.75), Inches(3.2), Inches(4.7), Inches(0.4),
     [{"text": "AFTER AUDIT + CLEANING", "size": 11, "color": GRAY}])
tbox(s, Inches(7.75), Inches(3.55), Inches(4.7), Inches(2.9),
     [{"runs": [("9,456", {"bold": True, "size": 20, "color": BLUE}),
                ("  modelling series rows", {"size": 13.5, "color": INK2})], "space_after": 8},
      {"runs": [("10,318", {"bold": True, "size": 20, "color": BLUE}),
                ("  modelling map rows", {"size": 13.5, "color": INK2})], "space_after": 8},
      {"runs": [("107", {"bold": True, "size": 20, "color": ORANGE}),
                ("  Cologne Major 2026 matches — ", {"size": 13.5, "color": INK2}),
                ("held completely outside development", {"size": 13.5, "color": ORANGE, "bold": True})],
       "space_after": 8, "spacing": 1.05},
      {"text": "Rejected rows are logged with reasons (missing format/score, ties) — never silently dropped.",
       "size": 12, "color": INK2, "spacing": 1.12}])
notes(s,
      "A large public dataset — but it needed a serious audit before it could be trusted; the winner label itself was broken.",
      "The data covers about nine and a half thousand professional series from January 2023 to late June 2026 — roughly ten and a half thousand played maps, "
      "across three tiers of tournaments. Two audit findings shaped everything. First, the dataset's own 'team1 wins' label disagreed with the actual match scores "
      "in half the rows — so the target was reconstructed directly from the scores. Second, team IDs turned out to be one-per-match, not persistent identities, so "
      "team identity was rebuilt from normalized names with a manual review of ambiguous cases. After conservative cleaning we keep about 9,450 series and 10,300 maps "
      "for modelling — and, crucially, the 107 Cologne Major matches were fenced off from day one, reserved as a purely external evaluation event.",
      ["Target reconstruction: winner = sign(score1 − score2); 5 ties and 1 missing score dropped.",
       "Every exclusion is logged with a reason — auditable, reproducible cleaning.",
       "106 of the 107 Cologne rows are official matches; 1 is a showmatch (explained on the Cologne slide)."],
      [("Where does the data come from?", "A public Kaggle export of professional CS2 matches (HLTV-style coverage), audited from scratch with reproducible pandas scripts."),
       ("How did you find the broken label?", "A raw-data audit compared the provided team1_win column against the actual scores — 49.9% disagreement, and internally inconsistent between series and map rows. Scores were self-consistent, so they define the target."),
       ("Is ~10k matches enough?", "Enough to learn a real signal above baselines, not enough for per-map or BO5 subtleties — that's in the limitations.")],
      "With clean data in hand, the pipeline itself is designed around one enemy: information leakage.")
footer(s, 4)

# ===========================================================================
# SLIDE 5 — METHODOLOGY
# ===========================================================================
s = new_slide()
header(s, 5, "Methodology", "A chronological pipeline built to prevent leakage")
steps = [("Raw match data", WHITE),
         ("Audit · cleaning · canonical team identity", WHITE),
         ("Chronological replay → historical state before every match", WHITE),
         ("Feature engineering (differences A − B)", WHITE),
         ("Chronological 70 / 15 / 15 split", WHITE),
         ("Tune + compare LR · RF · XGB (train-only CV)", WHITE),
         ("Freeze selected models", PANEL2),
         ("Sealed internal TEST + external Cologne evaluation", PANEL2),
         ("Application (API + PWA)", PANEL2)]
sx = ML
sy = Inches(1.5)
sw_ = Inches(3.86)
sh_ = Inches(0.52)
col_x = [ML, ML + Inches(4.12), ML + Inches(8.24)]
for i, (txt, f) in enumerate(steps):
    col = i // 3
    row = i % 3
    x = col_x[col]
    y = sy + row * Inches(0.74)
    pn = panel(s, x, y, sw_, sh_, fill=f, line=LINE)
    tf = pn.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, [{"runs": [(f"{i+1}. ", {"bold": True, "color": BLUE, "size": 12}),
                             (txt, {"size": 12.5, "bold": i >= 6})]}])
    if row < 2:
        arrow_d(s, x + Inches(1.83), y + sh_ + Inches(0.03), Inches(0.2), Inches(0.14), fill=LINE)
for col in range(2):
    arrow_r(s, col_x[col] + sw_ + Inches(0.02), sy + Inches(2.3), Inches(0.20), Inches(0.16), fill=LINE)

# split timeline
ty = Inches(4.05)
tbox(s, ML, ty, Inches(6.0), Inches(0.3),
     [{"text": "CHRONOLOGICAL SPLIT — the future is never in the training set", "size": 11.5,
       "color": GRAY}])
bar_y = ty + Inches(0.38)
bar_h = Inches(0.52)
seg = [(0.575, "TRAIN — 6,619 series", BLUE, WHITE, "Jan 2023 → Aug 2025"),
       (0.155, "VAL — 1,419", CYAN, WHITE, "→ Jan 2026"),
       (0.155, "TEST — 1,418 (sealed)", NAVY, WHITE, "→ Jun 2026"),
       (0.115, "COLOGNE", ORANGE, WHITE, "external")]
x = ML
total_w = CW
for frac, lab, fill, txtc, sub in seg:
    wseg = Emu(int(total_w * frac))
    r = rect(s, x, bar_y, wseg, bar_h, fill)
    tf = r.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, [{"text": lab, "size": 11.5, "bold": True, "color": txtc,
                    "align": PP_ALIGN.CENTER, "space_after": 0},
                   {"text": sub, "size": 9, "color": txtc, "align": PP_ALIGN.CENTER}])
    x += wseg
panel(s, ML, Inches(5.35), CW, Inches(1.5), fill=PANEL, line=BLUE)
tbox(s, ML + Inches(0.25), Inches(5.52), CW - Inches(0.5), Inches(1.2),
     [{"runs": [("Leakage prevention, concretely:  ", {"bold": True, "size": 14, "color": BLUE})],
       "space_after": 4},
      {"text": "every feature for a match is computed only from matches strictly earlier in time · matches sharing a timestamp never see each other's results · "
               "hyperparameters chosen by expanding-window CV inside TRAIN only · the TEST partition stayed sealed until both models were frozen · "
               "Cologne was never touched during development — training and live inference share the exact same feature code.",
       "size": 13, "color": INK2, "spacing": 1.18}])
notes(s,
      "The pipeline replays history in order, so every prediction uses only what was knowable at that moment — and evaluation gets strictly harder from left to right.",
      "The methodology is one long chronological pipeline. After the audit and cleaning, we replay all matches in time order, maintaining each team's evolving state — "
      "ratings, form, map history. When a match arrives, its features are computed from that state first, and only afterwards is its result applied. That single rule "
      "makes leakage structurally impossible rather than something we hope we avoided. The split is chronological too: train on 2023 to August 2025, validate on the "
      "following five months, and a sealed test partition covering spring 2026. Hyperparameters were tuned with expanding-window cross-validation inside the training "
      "period only. And on the far right, the hardest evaluation: the Cologne Major, completely outside development, evaluated only after the models were frozen.",
      ["Expanding-window CV = each fold trains on an earlier period and validates on the next one — CV that respects time.",
       "Timestamp batching: matches with identical timestamps get features first, results applied after — none sees another's outcome.",
       "Same build_features() function serves training rows and live predictions — train/inference code cannot diverge."],
      [("Why a chronological split instead of random?", "Random splits leak the future: a model would train on matches played after its validation matches, inflating scores. Deployment always predicts forward in time, so evaluation must too."),
       ("What is leakage exactly?", "Any information in a training feature that would not have been available at prediction time — e.g. using a team's season-end rating for a January match."),
       ("How do you know there is no leakage?", "By construction (strict-past feature computation, tested timestamp handling, sealed test, hashed frozen artifacts) — plus the external Cologne evaluation, which is immune to internal leakage by design.")],
      "Inside that pipeline, what do the features actually look like?")
footer(s, 5)

# ===========================================================================
# SLIDE 6 — FEATURES
# ===========================================================================
s = new_slide()
header(s, 6, "Feature engineering", "Team history, summarized into interpretable families")
fam = [
    ("Overall strength", "ELO rating maintained match-by-match (start 1500, K = 32) — elo_diff is the single strongest feature in both models.", BLUE),
    ("Recent form", "win rate over last 5/10 series · average score margins · time-decayed form", BLUE),
    ("Opponent strength", "average ELO of recent opponents — beating strong teams counts for more", BLUE),
    ("Map pool & map history", "pool depth, best-map strength, experience and record on the selected map", CYAN),
    ("Players & roster", "current players' recent performance (K/D, KAST) · roster continuity on the map", CYAN),
    ("Activity & experience", "matches in last 30 days · days since last match · total history volume", CYAN),
]
cw2 = Inches(3.94)
ch2 = Inches(1.52)
for i, (t, d, c) in enumerate(fam):
    x = ML + (i % 3) * (cw2 + Inches(0.14))
    y = Inches(1.5) + (i // 3) * (ch2 + Inches(0.18))
    panel(s, x, y, cw2, ch2, fill=PANEL)
    rect(s, x, y, Inches(0.07), ch2, c)
    tbox(s, x + Inches(0.22), y + Inches(0.12), cw2 - Inches(0.4), ch2 - Inches(0.2),
         [{"text": t, "size": 14, "bold": True, "space_after": 3},
          {"text": d, "size": 11.5, "color": INK2, "spacing": 1.1}])
panel(s, ML, Inches(5.15), CW, Inches(1.65), fill=WHITE, line=LINE)
tbox(s, ML + Inches(0.25), Inches(5.32), CW - Inches(0.5), Inches(1.35),
     [{"runs": [("Two design rules.  ", {"bold": True, "size": 14, "color": BLUE}),
                ("(1) Almost everything enters as a difference, Team A − Team B — the dataset's team ordering carried a spurious ~55% bias that a model could "
                 "otherwise exploit; differences plus mirrored training rows neutralize it.  (2) A team with no history gets explicit neutral defaults "
                 "(ELO 1500, rates 0.5) and the model is told history is missing — never a fake value.",
                 {"size": 13, "color": INK2})], "spacing": 1.18, "space_after": 6},
      {"runs": [("Scale:  ", {"bold": True, "size": 13}),
                ("17 features for the pre-veto series model  ·  131 for the known-map model.  No stat from the match being predicted is ever used.",
                 {"size": 13, "color": INK2})], "spacing": 1.1}])
notes(s,
      "Features are interpretable summaries of each team's past, expressed as A-minus-B differences — with ELO difference the strongest single signal.",
      "Every feature summarizes the past. Six families: overall strength — a classic ELO rating updated after every match; recent form over the last five and ten "
      "series; opponent strength, so beating strong opponents counts more than farming weak ones; the map families — how deep a team's map pool is and how good they "
      "are on the specific selected map; player and roster features — how the current five players have been performing; and activity — is the team rusty or "
      "match-sharp. Two design rules matter. Everything enters as a Team-A-minus-Team-B difference, because the raw data had a spurious ordering bias worth about "
      "five percentage points that a model would happily learn instead of real skill. And cold starts are explicit: an unknown team gets neutral defaults plus a "
      "flag saying 'no history', so the model can treat genuine uncertainty as such. The pre-veto model uses 17 of these; the known-map model 131.",
      ["ELO: expected score E = 1/(1+10^((R_B−R_A)/400)); update R ← R + 32·(result − E). Simple, interpretable, deliberately untuned.",
       "Mirrored augmentation: each training row also appears with teams swapped and the label flipped — the model can't learn 'team1 wins more'.",
       "elo_diff top feature in both models (RF permutation importance 0.071; XGB #1 as well)."],
      [("Why ELO and not something newer (Glicko/TrueSkill)?", "Transparent, one-line update, strong baseline; refinements (time decay, per-tier K) are listed as future work — the comparison to a pure-ELO baseline shows how much the ML adds on top."),
       ("131 features — overfitting risk?", "Controlled by chronological CV, strong regularization (final XGB uses depth-2 trees), and verified by the small train→test gap on the sealed test."),
       ("How are roster changes handled?", "Player-level features follow the current five players' own histories, and roster-continuity features measure how much of the map experience belongs to the current core.")],
      "Three model families competed on these features.")
footer(s, 6)

# ===========================================================================
# SLIDE 7 — MODELS COMPARED
# ===========================================================================
s = new_slide()
header(s, 7, "Models", "Three families, tuned identically, compared honestly")
data = [
    [("Model", {}), ("What it is", {}), ("Why it is here", {}), ("Character in this project", {})],
    [("Logistic Regression", {"bold": True}),
     ("Linear model on the feature differences — implemented from scratch (NumPy)", {}),
     ("Transparent baseline every tree model must beat", {}),
     ("Never overfits here (train ≈ validation); surprisingly competitive", {})],
    [("Random Forest", {"bold": True}),
     ("Many decorrelated decision trees vote; averaging tames variance", {}),
     ("Captures non-linear interactions with little assumption", {}),
     ("Untuned it memorized TRAIN (AUC gap +0.37); tuning cut the gap to +0.06", {})],
    [("XGBoost", {"bold": True}),
     ("Gradient-boosted trees — each tree corrects the previous ones' errors", {}),
     ("State of the art on tabular data; strong regularization", {}),
     ("Most stable train→validation transition (gap +0.007)", {})],
]
styled_table(s, ML, Inches(1.55), CW, Inches(3.3), data,
             col_widths=[Inches(2.2), Inches(3.6), Inches(3.1), Inches(3.19)],
             row_h=Inches(0.86), body_size=13)
panel(s, ML, Inches(5.3), CW, Inches(1.5), fill=PANEL)
tbox(s, ML + Inches(0.25), Inches(5.47), CW - Inches(0.5), Inches(1.2),
     [{"runs": [("Fair-comparison rules:  ", {"bold": True, "size": 14, "color": BLUE}),
                ("identical features, identical chronological split, identical train-only expanding-window CV folds for tuning each model, "
                 "identical mirrored-training policy — and each tuned model was scored exactly once on the held-out validation period. "
                 "Untuned-vs-tuned numbers are never mixed into an algorithm verdict.",
                 {"size": 13.5, "color": INK2})], "spacing": 1.2}])
notes(s,
      "Three model families with different inductive biases, given exactly the same information and tuning budget — so differences reflect the algorithms, not the setup.",
      "The three candidates. Logistic regression — a linear model, implemented from scratch — is the transparency baseline; if trees can't beat it, the extra "
      "complexity isn't earning its keep. Random forest builds hundreds of decision trees on random subsets and averages them; it captures non-linear structure, but "
      "untuned it simply memorized the training set — a train-validation AUC gap of plus 0.37 — and tuning its depth and leaf sizes brought that down to 0.06. "
      "XGBoost builds trees sequentially, each correcting the last, with heavy regularization — it had the most stable train-to-validation behaviour of the three. "
      "The comparison was kept deliberately fair: same features, same chronological folds for tuning, same augmentation, and one single scoring pass on validation "
      "per tuned model.",
      ["RF = bagging (parallel trees, variance reduction); XGB = boosting (sequential trees, bias reduction) — know this contrast.",
       "All three landed within ~0.015 AUC of each other → the features, not the algorithm, are the binding constraint.",
       "LR 'from scratch' = gradient-descent implementation, not sklearn."],
      [("What is a Random Forest?", "An ensemble of decision trees, each trained on a bootstrap sample with random feature subsets at each split; predictions are averaged. Averaging many overfit trees gives a low-variance ensemble."),
       ("What is XGBoost?", "Gradient boosting: trees added one at a time, each fit to the current errors (gradient of log loss), with shrinkage and regularization. Typically the strongest family on tabular data."),
       ("Why no neural networks?", "~10k rows of tabular data is where trees dominate; an MLP would add tuning cost and opacity with little expected gain. Listed as future work with more data.")],
      "So who won? Here are the validation results — and the selection logic.")
footer(s, 7)

# ===========================================================================
# SLIDE 8 — PRE-VETO RESULTS
# ===========================================================================
s = new_slide()
header(s, 8, "Results — pre-veto series model", "All three are close — probabilities decide the winner")
pic(s, PFIG / "preveto_model_comparison.png", ML, Inches(1.45), Inches(8.55), Inches(3.85),
    frame=False)
p = panel(s, Inches(9.45), Inches(1.55), Inches(3.25), Inches(3.65), fill=PANEL, line=BLUE)
tbox(s, Inches(9.68), Inches(1.72), Inches(2.82), Inches(3.4),
     [{"text": "SELECTION HIERARCHY (fixed in advance)", "size": 10.5, "color": GRAY, "space_after": 6},
      {"runs": [("1. Probability quality", {"bold": True, "size": 13}),
                ("  — Log Loss, Brier → RF", {"size": 12.5, "color": INK2})], "space_after": 5},
      {"runs": [("2. ROC-AUC", {"bold": True, "size": 13}),
                ("  → RF", {"size": 12.5, "color": INK2})], "space_after": 5},
      {"runs": [("3. Accuracy", {"bold": True, "size": 13}),
                ("  → LR (0.613)", {"size": 12.5, "color": INK2})], "space_after": 10},
      {"runs": [("Random Forest V2 selected", {"bold": True, "size": 14, "color": BLUE})],
       "space_after": 5},
      {"text": "The simulator consumes probabilities, not labels — so probability quality outranks raw accuracy.",
       "size": 12, "color": INK2, "spacing": 1.14}])
tbox(s, ML, Inches(5.55), CW, Inches(1.3),
     bullet_lines([
         {"runs": [("All three beat the majority baseline (0.553 accuracy) — and land within ~0.015 AUC of each other: the signal ceiling sits in the features.",
                    {"color": INK2})]},
         {"runs": [("LR has slightly higher accuracy (0.613 vs 0.607) — but RF V2 wins every probability metric, which is what the application actually uses.",
                    {"color": INK2})]},
         {"runs": [("Documented caveat: XGB V2's train→validation gap (+0.007) is far smaller than RF's (+0.055) — RF's win carries single-split risk.",
                    {"color": INK2})]},
     ], size=13, gap=6))
notes(s,
      "RF V2 was selected by a pre-registered hierarchy that puts probability quality first — even though LR had marginally higher accuracy.",
      "Here are the tuned models on the held-out validation period — 1,419 series the models never saw. Left panel: accuracy and AUC. Notice logistic regression "
      "actually has the best accuracy, by half a point. Right panel: probability quality — log loss and Brier, lower is better — and there random forest wins both, "
      "and also wins AUC. The selection rule was fixed before looking: probability quality first, then AUC, then accuracy — because everything downstream, "
      "especially the Monte-Carlo simulator, consumes probabilities, not thresholded labels. Under that hierarchy Random Forest V2 wins decisively. Two honest "
      "footnotes: all three models sit within one and a half AUC points of each other, which tells us the features, not the algorithm, are the constraint; and "
      "XGBoost had a much more stable train-to-validation transition, so RF's win on a single validation period carries some risk — that caveat is documented in "
      "the project, not hidden.",
      ["Numbers: LR/RF/XGB — Acc 0.613/0.607/0.612 · AUC 0.641/0.657/0.650 · LL 0.658/0.651/0.654 · Brier 0.233/0.230/0.231.",
       "Majority baseline 0.553 = always predicting the more frequent class.",
       "Selection hierarchy was written down before the comparison — no post-hoc metric shopping."],
      [("Why is AUC the right discrimination metric?", "AUC is threshold-free: the probability a random winner is ranked above a random loser. With a 55/45 class skew, accuracy alone is easy to game; AUC is not."),
       ("Isn't 61% accuracy low?", "Upsets are structural in CS2 — even bookmakers sit in the low-to-mid 60s. The value is in calibrated probabilities above baselines, not in certainty."),
       ("Why not ensemble the three?", "Considered and deliberately not pursued after freezing — the evaluation protocol forbids post-hoc model construction; listed as future work.")],
      "That was maps-unknown. When the maps are known, we can do better — with a different model.")
footer(s, 8)

# ===========================================================================
# SLIDE 9 — KNOWN-MAP MODEL
# ===========================================================================
s = new_slide()
header(s, 9, "Results — known-map model", "Predict each map, then compose the series exactly")
# left: DP diagram
lx = ML
panel(s, lx, Inches(1.5), Inches(5.6), Inches(2.5), fill=PANEL)
tbox(s, lx + Inches(0.2), Inches(1.62), Inches(5.2), Inches(0.3),
     [{"text": "BO3 EXAMPLE — XGB V3 per map, then dynamic programming", "size": 10.5, "color": GRAY}])
maps_ex = [("Mirage", "57%"), ("Inferno", "61%"), ("Nuke", "59%")]
for i, (m, pv) in enumerate(maps_ex):
    x = lx + Inches(0.25) + i * Inches(1.32)
    pn = panel(s, x, Inches(2.0), Inches(1.2), Inches(0.78), fill=WHITE, line=LINE)
    tf = pn.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, [{"text": m, "size": 11, "color": INK2, "align": PP_ALIGN.CENTER, "space_after": 1},
                   {"text": pv, "size": 15, "bold": True, "color": CYAN, "align": PP_ALIGN.CENTER}])
arrow_r(s, lx + Inches(4.28), Inches(2.28), Inches(0.3), Inches(0.2), fill=LINE)
pn = panel(s, lx + Inches(4.62), Inches(2.0), Inches(0.86), Inches(0.78), fill=WHITE, line=CYAN)
tf = pn.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_set_text(tf, [{"text": "DP", "size": 15, "bold": True, "color": CYAN, "align": PP_ALIGN.CENTER}])
tbox(s, lx + Inches(0.25), Inches(2.95), Inches(5.2), Inches(0.95),
     [{"runs": [("→  P(win series) = 66%", {"bold": True, "size": 16, "color": NAVY})], "space_after": 4},
      {"text": "State = (maps played, maps won). Map 3 only contributes where the series is still 1–1 — exact composition for any BO-N, no simulation needed.",
       "size": 11.5, "color": INK2, "spacing": 1.1}])
# left bottom: metrics chips
tbox(s, lx, Inches(4.28), Inches(5.6), Inches(0.3),
     [{"text": "SEALED INTERNAL TEST — 1,427 MAPS, MODEL FROZEN FIRST", "size": 10.5, "color": GRAY}])
mvals = [("0.613", "Accuracy"), ("0.649", "ROC-AUC"), ("0.652", "Log Loss"), ("0.230", "Brier")]
for i, (v, l) in enumerate(mvals):
    chip(s, lx + i * Inches(1.44), Inches(4.62), Inches(1.32), Inches(0.95), v, l,
         num_color=CYAN, num_size=18, label_size=10.5)
tbox(s, lx, Inches(5.75), Inches(5.75), Inches(1.1),
     [{"text": "95% bootstrap CI on AUC: 0.620 – 0.679 — clearly above chance. TEST slightly exceeded its development estimate (AUC 0.620) — no sign of overfitting to development history.",
       "size": 12, "color": INK2, "spacing": 1.15}])
# right: baselines figure + confusion matrix
pic(s, PFIG / "knownmap_test_vs_baselines.png", Inches(6.45), Inches(1.5), Inches(6.25), Inches(2.35), frame=False)
pic(s, FIG / "map_xgb_v3_test_confusion_matrix.png", Inches(7.3), Inches(4.05), Inches(4.4), Inches(2.85))
notes(s,
      "The known-map system predicts each map independently and composes the series probability exactly with dynamic programming — and it was validated once, on a sealed test.",
      "In known-maps mode, the XGBoost model scores each map separately — here Mirage 57, Inferno 61, Nuke 59 for Team A. A small dynamic program then walks "
      "through the series states — how many maps each side has won — and sums the probability of every path where Team A reaches two wins first. Map three only "
      "counts in branches where the series is one-one, so 'later maps only play if needed' is handled exactly, for any best-of-N. The map model was evaluated "
      "exactly once, on a sealed test of 1,427 maps closed until the model was frozen: accuracy 61.3 percent, AUC 0.649 with a bootstrap interval clearly above "
      "chance. The right chart gives context: it clearly beats a map-specific ELO baseline and a coin flip, and edges the overall-ELO baseline on probability "
      "quality; the confusion matrix shows reasonably balanced errors.",
      ["DP: dp[maps][wins] accumulates path probabilities; terminal when a side reaches ceil(N/2) wins.",
       "Composition assumes conditional independence of maps given the pre-match features — a stated simplification.",
       "'Team1 rate' in test is 57% — that's why recall on wins (0.71) is higher than on losses."],
      [("What is dynamic programming here?", "Exact expansion of the series over states (maps played, maps won): each map multiplies its win/lose probabilities into the surviving states; the answer is the total mass reaching the required wins. Not a heuristic, not sampling."),
       ("Are maps really independent?", "We assume independence given the features — momentum within a series is not modelled; a known limitation."),
       ("Why is the AUC here (0.649) lower than the series model's Cologne AUC?", "Different tasks — single maps are noisier than whole series; the numbers are not comparable.")],
      "Before the tournament part, one short slide on why probabilities are the bridge.")
footer(s, 9)

# ===========================================================================
# SLIDE 10 — WHY PROBABILITY MATTERS / BRIDGE TO MC
# ===========================================================================
s = new_slide()
header(s, 10, "From probabilities to tournaments", "One match probability → 50,000 possible Majors")
panel(s, ML, Inches(1.55), Inches(5.85), Inches(4.6), fill=PANEL)
tbox(s, ML + Inches(0.25), Inches(1.75), Inches(5.35), Inches(4.2),
     [{"text": "IF WE ONLY KEPT THE LABEL…", "size": 11, "color": GRAY, "space_after": 5},
      {"text": "“Vitality beats Falcons.”", "size": 17, "bold": True, "space_after": 6},
      {"text": "Then every simulated tournament is identical: the favorite wins every match, the top seed lifts the trophy, and an upset is literally impossible.",
       "size": 13.5, "color": INK2, "spacing": 1.18, "space_after": 14},
      {"text": "KEEPING THE PROBABILITY…", "size": 11, "color": BLUE, "space_after": 5},
      {"runs": [("winner ~ Bernoulli(p)", {"size": 19, "bold": True, "color": BLUE, "font": "Consolas"})],
       "space_after": 6},
      {"text": "Each simulated match is a weighted coin flip with the model's own probability: p = 0.62 → Team A wins 62% of simulated runs and loses 38% — upsets occur exactly as often as the model believes they should.",
       "size": 13.5, "color": INK2, "spacing": 1.18}])
panel(s, Inches(6.85), Inches(1.55), Inches(5.85), Inches(4.6), fill=WHITE, line=LINE)
tbox(s, Inches(7.1), Inches(1.75), Inches(5.35), Inches(4.2),
     [{"text": "REPEATED OVER A WHOLE BRACKET", "size": 11, "color": GRAY, "space_after": 8},
      {"runs": [("50,000", {"size": 30, "bold": True, "color": BLUE, "font": FONT_SB}),
                ("  complete simulated Majors", {"size": 15, "color": INK2})], "space_after": 8},
      {"text": "Every run: 106 matches sampled through the real Swiss + playoff rules, one full bracket at a time.",
       "size": 13.5, "color": INK2, "spacing": 1.16, "space_after": 12},
      {"runs": [("The output is a distribution, ", {"size": 14.5, "bold": True}),
                ("not a single guess:", {"size": 14.5, "color": INK2})], "space_after": 5},
      {"text": "“Vitality wins 30% of tournaments, Spirit 19%, NaVi 11%, Falcons 9%, …” — every team gets a championship, playoff and stage-exit probability with Monte-Carlo error bars.",
       "size": 13.5, "color": INK2, "spacing": 1.16, "space_after": 10},
      {"text": "This is why probability quality (Log Loss / Brier) — not accuracy — was the model selection criterion.",
       "size": 13, "color": BLUE, "bold": True, "spacing": 1.12}])
notes(s,
      "Sampling winners from the model's probabilities — winner ~ Bernoulli(p) — is what turns one match model into a tournament forecast with realistic upsets.",
      "This slide is the hinge of the talk. If we kept only labels, simulating a tournament would be pointless: the favorite wins every simulated match, so every "
      "simulation is the same bracket and upsets are impossible. Instead, each simulated match draws its winner like a weighted coin: if the model says 62 percent, "
      "Team A wins 62 percent of the simulated runs and loses the other 38. Crucially, we do not repeatedly pick whichever side has p above one half — we sample. "
      "Chaining those samples through the real tournament rules, 106 matches per run, and repeating fifty thousand times gives a distribution over entire "
      "tournaments: each team ends up with a probability of winning the Major, reaching playoffs, going out early. And that's precisely why the models were selected "
      "on probability quality — the simulator is only as honest as the probabilities you feed it.",
      ["Bernoulli(p) = single biased coin flip; the sampled winner advances and the bracket continues from that state.",
       "50,000 runs → Monte-Carlo standard error ≤ 0.22 percentage points on any probability — sampling noise is negligible.",
       "A deterministic favorite-path was also computed as a contrast — shown two slides ahead."],
      [("Why Monte Carlo instead of exact computation?", "The Swiss format's pairings depend on records and rematch-avoidance — the state space of full tournaments is combinatorially huge, so exact enumeration is infeasible; sampling approximates the distribution to arbitrary precision."),
       ("Why exactly 50,000?", "Chosen so Monte-Carlo noise (max SE ~0.22pp) is far below the differences we care about; it's also cheap because all 2,976 possible matchup probabilities are precomputed once."),
       ("Does the simulation update ratings between rounds?", "No — team state stays frozen at the pre-event snapshot throughout; simulated results never feed back into features.")],
      "Now the real test: we froze all of this before the Cologne Major and let reality grade it.")
footer(s, 10)

# ===========================================================================
# SLIDE 11 — COLOGNE EXTERNAL EVALUATION
# ===========================================================================
s = new_slide()
header(s, 11, "External evaluation", "IEM Cologne Major 2026 — frozen before, judged after")
tbox(s, ML, Inches(1.42), Inches(5.9), Inches(0.35),
     [{"text": "THE PROTOCOL", "size": 11.5, "color": GRAY}])
tbox(s, ML, Inches(1.78), Inches(5.9), Inches(3.3),
     bullet_lines([
         {"runs": [("Model, features and all team state frozen ", {"bold": True}),
                   ("before the event (cutoff 2 Jun 2026) — artifact hashes recorded.", {"color": INK2})]},
         {"runs": [("50,000 tournaments simulated pre-event", {"bold": True}),
                   (" and permanently saved — no Cologne result existed in the system.", {"color": INK2})]},
         {"runs": [("Only then were the real results opened ", {"bold": True}),
                   ("and compared — the frozen predictions were never regenerated.", {"color": INK2})]},
         {"runs": [("106 official matches scored ", {"bold": True}),
                   ("(one Germany–Poland showmatch excluded, with documented evidence).", {"color": INK2})]},
     ], size=13.5, gap=9))
panel(s, ML, Inches(5.15), Inches(5.9), Inches(1.6), fill=PANEL)
tbox(s, ML + Inches(0.22), Inches(5.32), Inches(5.5), Inches(1.3),
     [{"runs": [("Read this as a demonstration on one real event ", {"bold": True, "size": 13}),
                ("— every metric landed on the favorable side of its development-validation counterpart, but a single 106-match tournament is evidence, not statistical proof of improvement.",
                 {"size": 13, "color": INK2})], "spacing": 1.16}])
pic(s, PFIG / "cologne_match_metrics.png", Inches(6.7), Inches(1.55), Inches(6.05), Inches(3.35), frame=False)
mvals = [("0.642", "Accuracy (dev: 0.607)"), ("0.697", "ROC-AUC (dev: 0.657)"),
         ("0.632", "Log Loss (dev: 0.651)"), ("0.221", "Brier (dev: 0.230)")]
for i, (v, l) in enumerate(mvals):
    chip(s, Inches(6.72) + i * Inches(1.55), Inches(5.35), Inches(1.44), Inches(1.15), v, l,
         num_color=BLUE, num_size=19, label_size=9.5)
notes(s,
      "The strongest evidence in the project: a genuinely external, pre-registered evaluation — model and state frozen before the event, judged on 106 real Major matches.",
      "This is the evaluation I trust most, because it is immune to any internal mistake. Before the Major began, everything was frozen — the random forest, its "
      "preprocessing, and the team-state snapshot ending three days before the first match — with cryptographic hashes recorded. The full 50,000-run simulation was "
      "executed and saved before any result existed; only then were the real results opened. On the 106 official matches, the frozen model scored 64.2 percent "
      "accuracy, AUC 0.697, log loss 0.632 against 0.693 for an uninformed baseline, Brier 0.221 against 0.250. Every number is on the favorable side of the model's "
      "own development validation — encouraging, but I'll say precisely this much: a strong demonstration on one real event, not statistical proof. One data note: "
      "of 107 Cologne rows, one was a Germany-versus-Poland showmatch, excluded with documented evidence.",
      ["Frozen state = strict pre-Cologne snapshot (last update 30 May 2026); one team (THUNDERdOWNUNDER) was a genuine cold start.",
       "Dev-validation comparison: LL 0.632 vs 0.651, Brier 0.221 vs 0.230, AUC 0.697 vs 0.657, Acc 0.642 vs 0.607.",
       "Match orientation: team_a = better tournament seed at pairing time; baseline accuracy 0.528 = team_a prevalence."],
      [("Was Cologne in the training data?", "No. It was fenced off in the evaluation manifest from phase one, structurally absent from all feature tables, and the prediction state ends before the event. The deployed app's snapshot later includes Cologne, but the historical evaluation reads only the frozen artifacts."),
       ("What does AUC 0.697 mean?", "Take a random Cologne match the model got a winner and a loser for — with probability ~0.70 the model ranked the actual winner above the actual loser. 0.5 is chance."),
       ("Why not confidence intervals here?", "The same teams recur across the 106 matches, so an IID bootstrap would understate uncertainty; rather than report a too-narrow interval, the project reports point estimates and says 'single event' explicitly.")],
      "So what did those 50,000 simulated Majors actually predict — and what really happened?")
footer(s, 11)

# ===========================================================================
# SLIDE 12 — SIMULATION VS REALITY
# ===========================================================================
s = new_slide()
header(s, 12, "Simulation vs reality", "The champion was the model's #4 — visible, not chosen")
pic(s, PFIG / "cologne_champion_top8.png", ML, Inches(1.42), Inches(8.6), Inches(3.6), frame=False)
p = panel(s, Inches(9.5), Inches(1.5), Inches(3.2), Inches(3.45), fill=PANEL)
tbox(s, Inches(9.72), Inches(1.65), Inches(2.8), Inches(3.2),
     [{"text": "TEAM FALCONS — ACTUAL CHAMPION", "size": 10.5, "color": ORANGE, "space_after": 6},
      {"runs": [("8.9%", {"size": 26, "bold": True, "color": ORANGE, "font": FONT_SB}),
                ("  pre-event", {"size": 12.5, "color": INK2})], "space_after": 3},
      {"text": "rank 4 of 32 · ~3× the uniform 3.1% reference", "size": 12, "color": INK2,
       "spacing": 1.1, "space_after": 8},
      {"text": "Reach playoffs: 69%\nReach semifinal: 38%\nReach final: 19%",
       "size": 12.5, "color": NAVY, "spacing": 1.25, "space_after": 6},
      {"text": "A plausible contender in the model's own distribution — not a shock outside it.",
       "size": 11.5, "color": INK2, "spacing": 1.12}])
# got right / got wrong
gy = Inches(5.15)
panel(s, ML, gy, Inches(5.95), Inches(1.7), fill=WHITE, line=GREEN)
tbox(s, ML + Inches(0.2), gy + Inches(0.12), Inches(5.55), Inches(1.5),
     [{"text": "WHAT THE SYSTEM GOT RIGHT", "size": 10.5, "color": GREEN, "space_after": 4},
      {"text": "5 of 8 real playoff teams were in the pre-event top-8 · champion ranked top-4 · match accuracy strongest in Stage 1 (0.70) · most confident call (Spirit over 9z, p=0.79) was correct",
       "size": 12, "color": INK2, "spacing": 1.14}])
panel(s, Inches(6.77), gy, Inches(5.95), Inches(1.7), fill=WHITE, line=ORANGE)
tbox(s, Inches(6.97), gy + Inches(0.12), Inches(5.55), Inches(1.5),
     [{"text": "WHAT IT GOT WRONG", "size": 10.5, "color": ORANGE, "space_after": 4},
      {"text": "The 29.7% favorite (Vitality) did not win — and lost to 9z as a ~0.78 favorite, the biggest single miss · neither actual finalist was in the pre-event top-2 · playoff-stage accuracy only 4/7",
       "size": 12, "color": INK2, "spacing": 1.14}])
notes(s,
      "The Monte-Carlo distribution contained the real outcome — Falcons at 8.9%, rank 4 — while any single deterministic bracket structurally could not.",
      "Here is the frozen pre-event championship forecast, top eight of thirty-two. Vitality was the clear favorite at about 30 percent — and in orange, Team "
      "Falcons at 8.9 percent, rank four. Falcons won the Major. Two readings. Positively: 8.9 percent is nearly three times the uniform one-in-thirty-two "
      "reference, Falcons were top-four in the distribution, and five of the eight real playoff teams were in the model's top eight — what happened was a plausible "
      "draw from the forecast. Negatively, and I want to be equally clear: the favorite did not win, the single most confident wrong call was Vitality over 9z at "
      "78 percent, and neither actual finalist was in the pre-event top two. That contrast is the argument for distributions: a deterministic 'favorite always "
      "wins' bracket can only crown its own favorite — verified: it crowned Vitality — while the Monte-Carlo distribution carried the real champion with real weight.",
      ["Never claim the prediction was 'right' — claim the distribution was honest and informative.",
       "8.9% vs 3.1%: compare to the uniform 32-team reference, NOT to a coin flip.",
       "Deterministic favorite path: 6/8 correct Stage-1 advancers, but 0/2 finalists and 0/1 champion."],
      [("Isn't 8.9% just a miss?", "Under a proper scoring view, assigning 8.9% to the actual champion (when uniform gives 3.1%) is a better forecast than almost any single-bracket prediction; the point of a distribution is exactly that secondary contenders sometimes win."),
       ("Could any model have picked Falcons #1?", "Only by being badly overconfident elsewhere — the three teams above Falcons had 60% combined championship mass; pre-event evidence genuinely favored them."),
       ("How well were probabilities sized overall?", "Mean probability assigned to actual winners was 0.547 — modestly above chance, consistent with a competitive field; per-stage accuracy was best early (0.70 Stage 1) and hardest in playoffs (4/7).")],
      "A quick look inside the machine that generated those 50,000 tournaments.")
footer(s, 12)

# ===========================================================================
# SLIDE 13 — SIMULATION ENGINE
# ===========================================================================
s = new_slide()
header(s, 13, "Tournament engine", "The real Major format, replayed 50,000 times")
# Swiss stage diagram
sy0 = Inches(1.55)
stage_w = Inches(2.75)
stage_h = Inches(3.5)
stages = [("STAGE 1 — SWISS", "16 teams (seeds 17–32 region)", "BO1, BO3 at stakes"),
          ("STAGE 2 — SWISS", "8 advancers + 8 invited", "BO1, BO3 at stakes"),
          ("STAGE 3 — SWISS", "8 advancers + top-8 seeds", "all BO3")]
for i, (t, sub, fmt) in enumerate(stages):
    x = ML + i * Inches(3.06)
    panel(s, x, sy0, stage_w, stage_h, fill=PANEL)
    tbox(s, x + Inches(0.16), sy0 + Inches(0.12), stage_w - Inches(0.3), Inches(0.8),
         [{"text": t, "size": 12.5, "bold": True, "color": BLUE, "space_after": 1},
          {"text": sub, "size": 10.5, "color": INK2, "space_after": 1},
          {"text": fmt, "size": 10.5, "color": GRAY}])
    tbox(s, x + Inches(0.16), sy0 + Inches(1.1), stage_w - Inches(0.3), Inches(2.3),
         [{"text": "• same-record pairings\n   (1–0 vs 1–0, 2–1 vs 2–1 …)", "size": 10.5,
           "color": INK2, "spacing": 1.12, "space_after": 4},
          {"text": "• rematches avoided by rule", "size": 10.5, "color": INK2, "space_after": 4},
          {"runs": [("• 3 wins → advance (8 teams)", {"size": 10.5, "color": GREEN, "bold": True})],
           "space_after": 4},
          {"runs": [("• 3 losses → eliminated", {"size": 10.5, "color": ORANGE, "bold": True})],
           "space_after": 4},
          {"text": "• 33 matches per stage", "size": 10.5, "color": INK2}])
    if i < 2:
        arrow_r(s, x + stage_w + Inches(0.03), sy0 + Inches(1.6), Inches(0.25), Inches(0.22), fill=BLUE)
x = ML + 3 * Inches(3.06)
panel(s, x, sy0, Inches(2.90), stage_h, fill=WHITE, line=BLUE)
tbox(s, x + Inches(0.16), sy0 + Inches(0.12), Inches(2.6), Inches(3.3),
     [{"text": "PLAYOFFS", "size": 12.5, "bold": True, "color": BLUE, "space_after": 1},
      {"text": "8 teams, single elimination", "size": 10.5, "color": INK2, "space_after": 8},
      {"text": "Quarterfinals — BO3\nSemifinals — BO3", "size": 11, "color": INK2,
       "spacing": 1.2, "space_after": 4},
      {"runs": [("Grand Final — BO5", {"size": 11.5, "bold": True})], "space_after": 10},
      {"text": "7 matches\n→ 106 per tournament", "size": 10.5, "color": GRAY, "spacing": 1.15}])
arrow_r(s, x - Inches(0.28), sy0 + Inches(1.6), Inches(0.25), Inches(0.22), fill=BLUE)
panel(s, ML, Inches(5.35), CW, Inches(1.5), fill=PANEL, line=LINE)
tbox(s, ML + Inches(0.25), Inches(5.52), CW - Inches(0.5), Inches(1.25),
     [{"runs": [("Each simulated match:  ", {"bold": True, "size": 13.5}),
                ("winner ~ Bernoulli(p)", {"size": 13.5, "bold": True, "color": BLUE, "font": "Consolas"}),
                ("  from a precomputed 2,976-entry probability matrix (32 teams × 31 opponents × 3 formats, RF V2 pre-veto — future maps are unknown).  "
                 "Seeded RNG per run → every one of the 5.3 million simulated matches is exactly reproducible.  "
                 "Engine validity check: fed the real results, it reproduces the actual bracket 106/106.",
                 {"size": 13, "color": INK2})], "spacing": 1.2}])
notes(s,
      "The engine implements the real Valve Swiss + playoff rules exactly, samples every match from the frozen probability matrix, and is provably faithful — it reproduces the real tournament when given the real results.",
      "The Cologne format: three Swiss stages of sixteen teams. In a Swiss stage you always play someone with the same record — one-and-oh against one-and-oh — "
      "rematches avoided by rule; three wins advance, three losses eliminate, which takes exactly 33 matches per stage. Eight advance, eight join at the next "
      "stage, until the final eight play single-elimination playoffs — best-of-three until a best-of-five grand final; 106 matches per tournament. Every simulated "
      "match samples its winner from the frozen matrix of all 2,976 team-pair-format probabilities, computed once from the pre-veto model — future map picks don't "
      "exist. Seeded random streams make all 5.3 million simulated matches reproducible. And one detail I like: as a validity check the engine was fed the 106 real "
      "results — it reproduced the entire actual bracket, every pairing and stage transition, 106 out of 106.",
      ["Pairings within a stage depend on records AND rematch-avoidance — that's why simulation, not closed-form math.",
       "Probability matrix: 32×31 ordered pairs × 3 formats = 2,976 entries; the model is never called during simulation.",
       "Why RF V2 here: the tournament simulator predicts future matches whose map vetoes haven't happened."],
      [("How are the Swiss pairings decided exactly?", "By the Valve rulebook: seed-based first round (1v9…8v16), then same-record pools paired by current seeding with exhaustive rematch-minimizing search; a 'Difficulty Score' (Buchholz-style) orders teams mid-stage."),
       ("Why not use the known-map model in the simulator?", "Maps for future tournament matches are unknown — using it would require inventing vetoes. The pre-veto model is the honest tool for that job."),
       ("What varies between the 50,000 runs?", "Only the sampled winners; the probabilities, rules and seeds structure stay fixed. Each run's RNG stream is derived from a base seed plus the run index.")],
      "All of this is wrapped into an application you can actually use.")
footer(s, 13)

# ===========================================================================
# SLIDE 14 — APPLICATION
# ===========================================================================
s = new_slide()
header(s, 14, "Application", "The models, deployed — ML stays in Python")
# left: architecture stack
ax = ML
aw = Inches(4.35)
layers = [
    ("Frozen models", "RF V2 (pre-veto) · XGB V3 (known-map)", PANEL2),
    ("Python inference + explanation core", "same feature code as training · per-prediction factor attributions (TreeSHAP / tree-path)", WHITE),
    ("FastAPI  ·  /api/v1", "typed contracts · versioned snapshot · startup hash checks · Major simulation service", WHITE),
    ("Next.js PWA", "installable, offline-capable · /predict with both modes", WHITE),
]
ly = Inches(1.5)
for i, (t, d, f) in enumerate(layers):
    pn = panel(s, ax, ly, aw, Inches(0.98), fill=f, line=BLUE if i == 3 else LINE)
    tbox(s, ax + Inches(0.18), ly + Inches(0.10), aw - Inches(0.36), Inches(0.8),
         [{"text": t, "size": 13.5, "bold": True, "space_after": 2},
          {"text": d, "size": 10.5, "color": INK2, "spacing": 1.05}])
    if i < 3:
        arrow_d(s, ax + aw / 2 - Inches(0.09), ly + Inches(0.99), Inches(0.18), Inches(0.13), fill=BLUE)
    ly += Inches(1.13)
tbox(s, ax, Inches(6.15), aw + Inches(0.4), Inches(0.9),
     [{"text": "No model logic is duplicated in TypeScript — the frontend only renders API responses. "
               "Deployment snapshot: data through 28 Jun 2026 (includes Cologne); the historical Cologne replay stays frozen separately.",
       "size": 11, "color": INK2, "spacing": 1.14}])
# right: screenshots
shots = [
    (FIG / "phase10a" / "06_preveto_result_hero_1440.jpg",
     "/predict, maps-unknown (pre-veto) — series probability, favored team, model factors"),
    (FIG / "phase10a" / "11_map_breakdown_reach_leverage.jpg",
     "/predict, maps-known — per-map probability, chance each map is played, series leverage"),
]
sx0 = Inches(5.35)
sy1 = Inches(1.42)
for i, (p, cap) in enumerate(shots):
    pic(s, p, sx0, sy1, Inches(7.37), Inches(2.36), align="left")
    tbox(s, sx0, sy1 + Inches(2.40), Inches(7.37), Inches(0.24),
         [{"text": cap, "size": 10.5, "color": GRAY}])
    sy1 += Inches(2.78)
notes(s,
      "The whole system ships as a real product: Python owns every model decision, FastAPI exposes it with verified contracts, and a Next.js PWA renders it — including explanations.",
      "The application layers mirror the science. At the bottom, the two frozen models. Above them, a Python inference core reusing the exact same feature code as "
      "training — so the app cannot drift from the evaluation — plus an explanation core: every prediction ships with factor attributions from the model itself, "
      "TreeSHAP for XGBoost and an exact tree-path decomposition for the forest, phrased associationally, never causally. A FastAPI layer exposes this with typed "
      "contracts; at startup it re-verifies the hashes of every model and state file and refuses to serve if anything drifted. On top, a Next.js progressive web "
      "app: pick teams and format, choose maps-unknown or maps-known, and you get the probability bar, the per-map breakdown — including the chance each map is "
      "even played — and the 'why' factors. The Major simulator is exposed through the same API; its dedicated UI page is the next release step. One rule "
      "throughout: no ML logic is ever duplicated in the frontend.",
      ["Explanations: XGB uses exact TreeSHAP (built into XGBoost); RF uses Saabas-style tree-path attribution — exact for this forest, but not Shapley values.",
       "The deployed snapshot legitimately includes Cologne (data through 28 Jun 2026) — allowed because the historical evaluation was frozen first.",
       "App predictions are at the deployment snapshot cutoff — it is not live August-2026 data."],
      [("Why FastAPI?", "Thin, typed, async Python layer — the models are already Python; pydantic contracts give validation for free; zero model logic in the transport layer."),
       ("Frontend predictions match the research code exactly?", "Yes — the API calls the same frozen pipelines; startup hash checks plus golden-fixture tests pin the numbers."),
       ("Are the explanations causal?", "No, and the UI says so — they describe what moved this model's prediction, not why a team will actually win.")],
      "To close: what worked, what the limits are, and where this goes next.")
footer(s, 14)

# ===========================================================================
# SLIDE 15 — CONCLUSION
# ===========================================================================
s = new_slide()
header(s, 15, "Conclusion", "What worked, what is limited, what comes next")
colw = Inches(3.94)
cols = [
    ("WHAT WORKED", GREEN, [
        "Real predictive signal above every baseline — on a sealed test and on an external Major",
        "Probabilities, not labels — selected and scored with proper scoring rules",
        "A frozen, hash-verified pre-event forecast honestly compared against reality",
        "A complete deployed system: engine → API → PWA, with per-prediction explanations",
    ]),
    ("LIMITATIONS", ORANGE, [
        "Data ends 28 Jun 2026 — the app predicts from that snapshot, not live form",
        "One external event; BO5 nearly absent (n=1 at Cologne); some maps thin",
        "Simple ELO (no decay/tier weighting) · maps assumed independent within a series",
        "Roster changes only partially modelled; team identity relies on name resolution",
    ]),
    ("NEXT", BLUE, [
        "Second data source (e.g. GRID) for freshness and cross-source identity",
        "Time-decayed form and opponent-adjusted rating refinements",
        "Deeper roster/player modelling (transfers, role changes)",
        "Probability calibration study + ensembling of the model families",
    ]),
]
for i, (t, c, items) in enumerate(cols):
    x = ML + i * (colw + Inches(0.14))
    panel(s, x, Inches(1.5), colw, Inches(4.35), fill=PANEL)
    rect(s, x, Inches(1.5), colw, Inches(0.09), c)
    tbox(s, x + Inches(0.2), Inches(1.72), colw - Inches(0.4), Inches(0.35),
         [{"text": t, "size": 13, "bold": True, "color": c}])
    tbox(s, x + Inches(0.2), Inches(2.14), colw - Inches(0.4), Inches(3.6),
         bullet_lines([{"runs": [(it, {"color": INK2})]} for it in items],
                      size=11.8, gap=9, marker_color=c))
panel(s, ML, Inches(6.05), CW, Inches(0.85), fill=WHITE, line=BLUE)
tbox(s, ML + Inches(0.25), Inches(6.2), CW - Inches(0.5), Inches(0.6),
     [{"runs": [("Takeaway:  ", {"bold": True, "size": 14, "color": BLUE}),
                ("pre-match information carries real, honestly-measurable signal about CS2 outcomes — and a probabilistic system can say so while telling you exactly how uncertain it is.  Thank you — questions welcome.",
                 {"size": 14, "color": NAVY})], "spacing": 1.15}])
notes(s,
      "The project delivers a modest but real, honestly measured predictive signal, end-to-end from raw data to a deployed probabilistic product — with its limits stated as clearly as its wins.",
      "To conclude. What worked: genuine predictive signal above every baseline, demonstrated twice — on a sealed internal test and on a fully external Major with "
      "everything frozen in advance; a system that is probabilistic end to end, selected and judged with proper scoring rules; and a complete shipped application. "
      "The limitations are equally explicit: the data ends June 28th 2026, so the app reasons from that snapshot; Cologne is one event and best-of-five is nearly "
      "absent; the ELO is deliberately simple; maps are treated as independent within a series; roster dynamics are only partially captured. Next steps follow "
      "directly: a second data source such as GRID, time-decayed and opponent-adjusted ratings, deeper roster modelling, and a proper calibration and ensembling "
      "study. The takeaway in one sentence: pre-match information carries real, honestly measurable signal about CS2 outcomes — and a probabilistic system can "
      "deliver it while being exact about its own uncertainty. Thank you.",
      ["Keep the three columns in this order if asked to summarize: worked → limited → next.",
       "If pressed for the single biggest limitation: data freshness (snapshot ends 28 Jun 2026) plus single-event external evidence."],
      [("What would you do first with more time?", "Ingest a second, fresher data source — it simultaneously fixes staleness, strengthens team identity, and enables a second external evaluation event."),
       ("Could this generalize to other esports?", "The architecture (chronological state engine → probabilistic model → DP/Monte-Carlo composition) is game-agnostic; only the feature families are CS2-specific."),
       ("Is the model good enough to bet on?", "That was never the goal and no market comparison was done; the honest claim is 'meaningfully better than uninformed baselines, with calibrated-looking probabilities'.")],
      "— end of talk —")
footer(s, 15)

prs.save(OUT)
print("saved", OUT)
