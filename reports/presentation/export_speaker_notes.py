# Dumps the embedded speaker notes of the deck into a markdown file,
# so the .md notes can never drift from the .pptx.
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
PPTX = Path(__file__).resolve().parent / "cs2_match_prediction_presentation.pptx"
OUT = Path(__file__).resolve().parent / "presentation_speaker_notes.md"

TITLES = [
    "Title — Predicting CS2 Matches with Probabilities",
    "Objective — Estimate P(Team A wins) before the match starts",
    "System overview — One question, two prediction modes",
    "Data — 3.5 years of professional CS2 matches",
    "Methodology — A chronological pipeline built to prevent leakage",
    "Feature engineering — Team history, summarized into interpretable families",
    "Models — Three families, tuned identically, compared honestly",
    "Results (pre-veto) — All three are close; probabilities decide the winner",
    "Results (known-map) — Predict each map, then compose the series exactly",
    "From probabilities to tournaments — One match probability -> 50,000 possible Majors",
    "External evaluation — IEM Cologne Major 2026, frozen before, judged after",
    "Simulation vs reality — The champion was the model's #4",
    "Tournament engine — The real Major format, replayed 50,000 times",
    "Application — The models, deployed; ML stays in Python",
    "Conclusion — What worked, what is limited, what comes next",
]

prs = Presentation(str(PPTX))
lines = [
    "# Speaker Notes — CS2 Match Prediction (15-minute presentation)",
    "",
    "Auto-extracted from `reports/presentation/cs2_match_prediction_presentation.pptx` by",
    "`reports/presentation/export_speaker_notes.py` — the same notes are embedded in the",
    "PPTX notes pane. Target pace: ~50-60 seconds per slide => ~13.5 minutes total,",
    "leaving safety margin inside the 15-minute limit.",
    "",
]
for i, slide in enumerate(prs.slides, start=1):
    lines.append(f"## Slide {i} — {TITLES[i-1]}")
    lines.append("")
    txt = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
    for raw in txt.splitlines():
        if raw in ("MAIN MESSAGE", "SCRIPT (speak naturally)", "WHAT I MUST UNDERSTAND",
                   "LIKELY PROFESSOR QUESTIONS", "TRANSITION"):
            lines.append(f"**{raw}**")
            lines.append("")
        elif raw.startswith("Q: "):
            lines.append(f"- *{raw}*")
        elif raw.startswith("A: "):
            lines.append(f"  - {raw}")
        elif raw.startswith("- "):
            lines.append(raw)
        elif raw.strip() == "":
            lines.append("")
        else:
            lines.append(raw)
    lines.append("")
OUT.write_text("\n".join(lines), encoding="utf-8")
print("wrote", OUT)
